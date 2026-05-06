"""Build the Authrex MVP pitch deck from the Cognizant Technoverse template.

Loads the official Technoverse_MVP_Presentation Template.pptx, fills the cover
+ 5 native content slides + inserts 5 additional content slides (UVP, Business
Model, ROI+Market+GTM, Agentic Workflow, Pilot Ask), preserves the closing +
eval criteria slides. Output: ops/demo/AUTHREX_MVP_DECK.pptx.

Run:
    cd backend && .venv/Scripts/python.exe -m scripts.build_deck
"""
from __future__ import annotations

import copy
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


AUTHREX_ROOT = Path(__file__).resolve().parents[2]   # .../Authrex
WORKSPACE_ROOT = AUTHREX_ROOT.parent                 # .../cts-technoverse
TEMPLATE_PATH = WORKSPACE_ROOT / "official-docs" / "mvp-ppt-template" / "Technoverse_MVP_Presentation Template.pptx"
OUT_PATH = AUTHREX_ROOT / "ops" / "demo" / "AUTHREX_MVP_DECK.pptx"


def _set_text(shape, text: str, *, font_size: int | None = None, bold: bool = False, color: tuple[int, int, int] | None = None) -> None:
    """Replace all text in shape with `text`. Preserves first paragraph's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def _set_bullets(shape, bullets: list[str | tuple[str, int]], *, font_size: int = 12, color: tuple[int, int, int] = (40, 40, 40)) -> None:
    """Replace text with bulleted list. Each item is text or (text, indent_level).
    Indent level 0 = top-level bullet; 1 = sub-bullet."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(*color)


def _find_shape(slide, name_or_text: str):
    """Find shape by name OR by current text content (substring)."""
    for shape in slide.shapes:
        if shape.name == name_or_text:
            return shape
        if shape.has_text_frame and name_or_text in shape.text_frame.text:
            return shape
    return None


def _duplicate_slide(prs: Presentation, source_idx: int) -> int:
    """Duplicate slide at source_idx, returning the new slide's index.

    python-pptx doesn't have a built-in slide-duplicate API; we manipulate the
    underlying XML to clone the slide and append it to the deck.
    """
    from copy import deepcopy
    from pptx.oxml.ns import qn

    source = prs.slides[source_idx]
    # Create a new blank slide using the same layout
    new_slide = prs.slides.add_slide(source.slide_layout)

    # Remove default placeholders from the new slide (they'll be re-created from copy)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy each shape's XML from source to new
    for shape in source.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)

    return len(prs.slides) - 1


def _move_slide_to(prs: Presentation, from_idx: int, to_idx: int) -> None:
    """Move slide from from_idx to to_idx in the slide order."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    sld = slides[from_idx]
    xml_slides.remove(sld)
    xml_slides.insert(to_idx, sld)


# =============================================================================
# Content for every slide
# =============================================================================


COVER = {
    "idea_title": "Authrex",
    "subtitle": "Oncology Prior-Authorization Copilot — Cognizant TriZetto AI Gateway-native specialty agent bundle",
    "team_lines": [
        "Idea Title : Authrex — close the AI velocity gap for oncology prior auth",
        "Team Captain Name : Preethi Sivachandran",
        "Team Members : <member 2> · <member 3> · <member 4>",
        "College Name : <college>",
    ],
}


# Native Slide 2 — Problem Statement | Solution Description
SLIDE_2 = {
    "left_section": "Problem Statement",
    "left_bullets": [
        "Cancer can't wait. Prior auth does.",
        ("18 min/case manual baseline · 7-day SLA · $35B/yr US PA admin cost", 1),
        ("36% of oncologists report a patient death linked to PA delay (ASCO)", 1),
        ("CMS-0057-F live Jan 1 2026; 95% of GenAI pilots fail to scale (MIT NANDA)", 1),
        "Stakeholders: provider care coordinators, oncologists, payer reviewers, members.",
    ],
    "right_section": "Solution Description",
    "right_bullets": [
        "Authrex — FHIR bundle to TriZetto-native decision in 90 seconds.",
        ("7-agent LangGraph DAG · 22 sub-agents · Bedrock + Claude Sonnet 4.6 · MCP-native", 1),
        ("HITL gate at confidence < 0.75 (CA SB 1120 / CMS § IV.C compliant)", 1),
        ("Submits as Facets v3 + QNXT v2 events to Cognizant TriZetto AI Gateway", 1),
        ("Per-case Evidence Pack with SHA-256 tamper hash — auditor-grade", 1),
    ],
}


# Native Slide 3 — Uniqueness / Innovativeness | Business Impact
SLIDE_3 = {
    "left_section": "Uniqueness / Innovativeness",
    "left_bullets": [
        "Not a chatbot. A CMS-0057-F-evidenced agent bundle.",
        ("First specialty agent bundle for Cognizant TriZetto AI Gateway (Aug 6, 2025 launch)", 1),
        ("Same Bedrock + Claude + MCP stack Cognizant standardized on Nov 4, 2025", 1),
        ("Cognizant Neuro-SAN compatible — drop-in HOCON network in ops/cognizant-neuro/", 1),
        ("Five enforcement layers: in-process Gateway → AWS API Gateway → VPC endpoint → IAM → Guardrails", 1),
        ("8 ADRs in canonical Nygard format + 18 architecture docs + live /architecture/layers introspection", 1),
    ],
    "right_section": "Business Impact",
    "right_bullets": [
        "$1,499.55 saved per case. $1.26B per half-star at Humana scale.",
        ("Cycle-time: 18 min → 90s (≈98% reduction; upper band of 2026 enterprise GenAI benchmarks)", 1),
        ("Productivity uplift: 38% (inside 20-40% benchmark band for AI copilots)", 1),
        ("Escalation reduction: 27% (inside 20-30% benchmark band)", 1),
        ("Star Ratings projection: +0.2 to +0.4 stars on PA-influenced measures", 1),
        ("Provider abrasion: ~25 minutes/case returned to clinical practice", 1),
    ],
}


# Native Slide 4 — Technical Design and architecture (full-width)
SLIDE_4 = {
    "section": "Technical Design and architecture",
    "bullets": [
        "5 named layers, live introspectable at GET /api/v1/architecture/layers",
        ("Experience Layer — React 18 SPA · 17 routes · SSE trace stream · role-aware (coordinator/reviewer/admin)", 1),
        ("Orchestration & Policy Engine — FastAPI + LangGraph 7-agent DAG · BudgetTracker · review_gate HITL · case_jobs queue (Postgres SKIP LOCKED)", 1),
        ("Context Retrieval Service — pluggable: Bedrock KB / Amazon Q Business / S3 Vectors substrate (one Pydantic schema)", 1),
        ("GenAI Gateway — literal app/llm/gateway.py · per-tenant model allowlist · 24h rolling token+USD quota · audit row per call", 1),
        ("Telemetry & Governance — Prometheus /metrics · 7 SLOs with PagerDuty burn-rate alerts · Evidence Pack with SHA-256 · Responsible AI model card (NIST AI RMF + ISO 42001 + EU AI Act)", 1),
        "AWS Foundation: Bedrock + Sonnet 4.6 + Haiku 4.5 · Bedrock KB · Bedrock Guardrails (per-tenant) · AgentCore Runtime (apply-ready) · Amazon Q Business · RDS Aurora Global · KMS multi-region · IRSA · NetworkPolicy",
    ],
}


# Native Slide 5 — Scalable / Reusable | Roadmap
SLIDE_5 = {
    "left_section": "Scalable / Reusable",
    "left_bullets": [
        "Multi-tenant by construction (organization_id boundary on every row)",
        ("4 apply-ready Terraform modules: multi-region · provisioned-throughput · bedrock-vpc-endpoint · s3-vectors", 1),
        ("Capacity model: 1K → 10K → 100K cases/day; per-tier sizing in ops/SCALING.md", 1),
        ("Customer onboarding: 7-15 business days end-to-end (per-tenant Bedrock Guardrail + KMS + TriZetto Gateway URL)", 1),
        ("New specialty (cardiology, behavioral health) = 3 markdown files + Kiro Hook regen", 1),
        ("CI/CD: .github/workflows/ — OIDC · Inspector · Semgrep · CycloneDX SBOM · canary deploy with auto-promote on error budget", 1),
    ],
    "right_section": "Roadmap",
    "right_bullets": [
        "Day 0 — TriZetto pilot customer signs (one Cognizant Facets payer)",
        ("Day 7 — per-tenant Bedrock Guardrail provisioned · TriZetto Gateway URL configured", 1),
        ("Day 21 — first production case live with full evidence pack", 1),
        ("Day 45 — Bedrock Provisioned Throughput pinned (1 MU Sonnet OneMonth commit)", 1),
        ("Day 60 — second specialty (cardiology) live via Kiro spec edit", 1),
        ("Day 90 — first pilot ROI report published (joint AWS + Cognizant blog post)", 1),
    ],
}


# Additional slides (inserted between slide 5 and slide 6)
ADDITIONAL_SLIDES = [
    {
        "section": "Unique Value Proposition",
        "bullets": [
            "Drop into TriZetto Monday. Audit-grade by Day 21. Star lift by Day 90.",
            ("Same Bedrock + Claude Sonnet 4.6 + MCP stack Cognizant standardized on (Nov 4, 2025 partnership)", 1),
            ("Day-1 add-on inside an existing Cognizant TriZetto subscription — no new platform decision", 1),
            ("CMS-0057-F evidence shipped as a first-class endpoint, not a slide claim", 1),
            ("Live /api/v1/compliance/case/{id} returns 8-clause scorecard with evidence pointers", 1),
            ("Live /api/v1/cases/{id}/evidence-pack returns single SHA-256-hashed JSON bundle (auditor-grade)", 1),
        ],
    },
    {
        "section": "Business Model & Cognizant Channel",
        "bullets": [
            "$5/case license. 300× customer headroom. Cognizant TriZetto sales channel.",
            ("Vs $1,500 AMA-loaded manual cost = $5,475K/yr displaced cost per 10K-case/day customer", 1),
            ("Bundled into existing TriZetto subscription as a 'Specialty Agent Bundle' SKU", 1),
            ("Standard Cognizant rev-share; same motion as any AI Gateway agent listing", 1),
            ("Marketed via Cognizant Health Sciences vertical (~30.1% of Cognizant FY2024 revenue)", 1),
            ("Addressable Day-0: ~80M Facets lives + ~20M QNXT lives ≈ 100M lives", 1),
        ],
    },
    {
        "section": "Market, ROI Calculator & GTM",
        "bullets": [
            "100M lives addressable. Live /roi calculator with Humana 6M slider preset.",
            ("Star Ratings revenue lift: $2.1M / 10K members / 0.5 stars (Lilac 2025)", 1),
            ("At Humana scale (~6M MA enrollees): $1.26B per half-star", 1),
            ("AHIP/BCBSA Apr 24 2026 commitment: 50 insurers signed for FHIR PA APIs", 1),
            ("Insurers eliminated only 11% of PAs 6 months in (the gap Authrex closes)", 1),
            ("Cognizant ask: 1 Facets pilot · 30 days · joint AWS+Cognizant blog · AWS Marketplace listing under Cognizant Bedrock Healthcare seller account", 1),
        ],
    },
    {
        "section": "Agentic Workflow & Why This Stack",
        "bullets": [
            "User goal → 7-agent network → 5 typed actions → auditable outcome.",
            ("Closes Cognizant's 'AI velocity gap' (Ravi Kumar Dec 2025: $500B AI infra spent, value missing)", 1),
            ("Closes the 'AI adaptation gap' — embeds into existing TriZetto workflow, doesn't replace it", 1),
            ("Aligned with Cognizant Flowsource UX shape — async submit, operator micro-steers via review_gate", 1),
            ("Bedrock AgentCore Action Groups: persist_decision · route_to_review · submit_to_trizetto_gateway · draft_appeal · notify_patient", 1),
            ("Anthropic partnership Nov 4, 2025 (Claude across 350K Cognizant employees) — same stack, same vendor", 1),
        ],
    },
    {
        "section": "The Pilot Ask",
        "bullets": [
            "One Cognizant Facets customer. 30 days. We close the AI velocity gap together.",
            ("TriZetto product team review of our Facets v3 + QNXT v2 schemas (built from public docs)", 1),
            ("Pilot kickoff: AeroFyta supplies engineering; Cognizant supplies the payer relationship", 1),
            ("Joint AWS + Cognizant blog post (same model as the re:Invent 2025 IND210 collaboration)", 1),
            ("AWS Marketplace listing under the Cognizant Bedrock Healthcare seller account", 1),
            ("Day 0 → Day 90 plan in ops/demo/COGNIZANT_GO_TO_MARKET.md — every milestone owner-tagged", 1),
        ],
    },
]


# =============================================================================
# Build the deck
# =============================================================================


def build_deck() -> Path:
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE_PATH}")

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE_PATH))

    n_slides_template = len(prs.slides)
    print(f"Loaded template: {n_slides_template} slides")

    # ---- Slide 1: cover ----
    slide1 = prs.slides[0]
    title = _find_shape(slide1, "MVP Presentation")
    if title is not None:
        _set_text(title, "Authrex", font_size=36, bold=True)
    text_block = _find_shape(slide1, "Idea Title")
    if text_block is not None:
        _set_bullets(text_block, COVER["team_lines"], font_size=14)
    print("  slide 1 (cover)              filled")

    # ---- Slide 2: Problem Statement | Solution Description ----
    slide2 = prs.slides[1]
    # Title at top stays as <<Idea Title>> placeholder — replace with "Authrex"
    title2 = _find_shape(slide2, "<<Idea Title>>")
    if title2 is not None:
        _set_text(title2, "Authrex", font_size=14, bold=True)
    # Find the two content placeholders by their existing default text
    left_content = _find_shape(slide2, "Explain the problem")
    if left_content is not None:
        _set_bullets(left_content, SLIDE_2["left_bullets"], font_size=11)
    right_content = _find_shape(slide2, "addressing the problem")
    if right_content is not None:
        _set_bullets(right_content, SLIDE_2["right_bullets"], font_size=11)
    print("  slide 2 (Problem | Solution) filled")

    # ---- Slide 3: Uniqueness | Business Impact ----
    slide3 = prs.slides[2]
    title3 = _find_shape(slide3, "<<Idea Title>>")
    if title3 is not None:
        _set_text(title3, "Authrex", font_size=14, bold=True)
    left3 = _find_shape(slide3, "Aspects of the solution")
    if left3 is not None:
        _set_bullets(left3, SLIDE_3["left_bullets"], font_size=11)
    right3 = _find_shape(slide3, "Outcome of implementing")
    if right3 is not None:
        _set_bullets(right3, SLIDE_3["right_bullets"], font_size=11)
    print("  slide 3 (Uniqueness | Impact) filled")

    # ---- Slide 4: Technical Design and architecture (full-width) ----
    slide4 = prs.slides[3]
    title4 = _find_shape(slide4, "<<Idea Title>>")
    if title4 is not None:
        _set_text(title4, "Authrex", font_size=14, bold=True)
    tech = _find_shape(slide4, "List the different technologies")
    if tech is not None:
        _set_bullets(tech, SLIDE_4["bullets"], font_size=11)
    print("  slide 4 (Technical Design)    filled")

    # ---- Slide 5: Scalable / Reusable | Roadmap ----
    slide5 = prs.slides[4]
    title5 = _find_shape(slide5, "<<Idea Title>>")
    if title5 is not None:
        _set_text(title5, "Authrex", font_size=14, bold=True)
    left5 = _find_shape(slide5, "Scope for the solution")
    if left5 is not None:
        _set_bullets(left5, SLIDE_5["left_bullets"], font_size=11)
    right5 = _find_shape(slide5, "future potential roadmap")
    if right5 is not None:
        _set_bullets(right5, SLIDE_5["right_bullets"], font_size=11)
    print("  slide 5 (Scalable | Roadmap)  filled")

    # ---- Insert 5 ADDITIONAL slides between slide 5 and slide 6 ----
    # We duplicate slide 4 (full-width single section) for each additional slide,
    # because additional slides have one Pentagon-arrow header (matches user's
    # "1 header per slide" intent).
    for i, extra in enumerate(ADDITIONAL_SLIDES):
        # Duplicate slide 4 (the full-width layout) → new slide is appended
        new_idx = _duplicate_slide(prs, 3)
        new_slide = prs.slides[new_idx]
        # Update header (Pentagon arrow currently says "Technical Design and architecture")
        header = _find_shape(new_slide, "Technical Design and architecture")
        if header is not None:
            _set_text(header, extra["section"], font_size=18, bold=True, color=(255, 255, 255))
        # Update title (top-left)
        title_x = _find_shape(new_slide, "<<Idea Title>>") or _find_shape(new_slide, "Authrex")
        if title_x is not None:
            _set_text(title_x, "Authrex", font_size=14, bold=True)
        # Update body content
        body = _find_shape(new_slide, "List the different technologies") or _find_shape(new_slide, "5 named layers")
        if body is not None:
            _set_bullets(body, extra["bullets"], font_size=11)
        # Move from end-of-deck to position just after slide 5 (idx 4 → 5+i)
        target_idx = 5 + i
        _move_slide_to(prs, len(prs.slides) - 1, target_idx)
        print(f"  added slide {target_idx + 1} ({extra['section']!s})")

    # Update slide 6 (Additional slides) — repurpose as the AeroFyta journey
    journey_idx = 5 + len(ADDITIONAL_SLIDES)
    journey_slide = prs.slides[journey_idx]
    journey_subtitle = _find_shape(journey_slide, "Narrative of your Technoverse journey")
    if journey_subtitle is not None:
        _set_text(
            journey_subtitle,
            "From idea (Apr 4) to MVP (May 6 Pune) to first prize (May 7) — Team AeroFyta's 33-day sprint. Built 7 parents · 22 sub-agents · 56 routes · 18 architecture docs · 8 ADRs · 4 Terraform modules · 1 Cognizant Impact Pack.",
            font_size=14,
        )
    print(f"  slide {journey_idx + 1} (journey)              updated")

    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    return OUT_PATH


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")  # type: ignore[attr-defined]
    out = build_deck()
    if not out.exists():
        sys.exit(1)
    print(f"OK: {out.stat().st_size} bytes")
